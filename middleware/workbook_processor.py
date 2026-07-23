"""
workbook_processor.py — SAP SuccessFactors Configuration Workbook Processor

Handles:
  1. Country extraction from SOW / SP51 documents via Claude AI
  2. Loading grounded reference workbooks (up to 4 slots)
  3. Identifying CSF (Country-Specific Feature) vs Global sheets
  4. Filtering CSF sheets to in-scope countries only
  5. Generating a filtered Configuration Workbook as .xlsx

SAP SF Employee Central knowledge applied:
  - CSF sheets are country-specific: named with ISO codes, contain "CSF" in tab name,
    or have a Country/Country Code column as a primary key column
  - Global sheets apply to all countries (Foundation Objects, Org Structure, Pay Components, etc.)
  - Common CSF indicators: tab names like DEU, GBR, USA, SGP, AUS, or prefixes like
    EC_CSF_, Pay_Info_, Legal_Entity_, Tax_ followed by a country code
  - A "Country" or "Country Group" column is the definitive CSF row filter
"""

import io
import os
import re
import json
import base64
import struct
import asyncio
from pathlib import Path
from collections import OrderedDict
from typing import Dict, List, Optional, Tuple, Any

import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter

import anthropic

# ── ISO 3166-1 country reference (alpha-2, alpha-3, primary name) ────────────
# Full ISO country list — SAP SF Employee Central ships 100+ localizations, so
# detection must not be limited to a hand-picked subset (that previously
# caused most in-scope countries to be silently dropped during workbook scan).
_ISO_COUNTRIES: List[Tuple[str, str, str]] = [
    ("AF", "AFG", "Afghanistan"), ("AL", "ALB", "Albania"), ("DZ", "DZA", "Algeria"),
    ("AS", "ASM", "American Samoa"), ("AD", "AND", "Andorra"), ("AO", "AGO", "Angola"),
    ("AI", "AIA", "Anguilla"), ("AG", "ATG", "Antigua and Barbuda"), ("AR", "ARG", "Argentina"),
    ("AM", "ARM", "Armenia"), ("AW", "ABW", "Aruba"), ("AU", "AUS", "Australia"),
    ("AT", "AUT", "Austria"), ("AZ", "AZE", "Azerbaijan"), ("BS", "BHS", "Bahamas"),
    ("BH", "BHR", "Bahrain"), ("BD", "BGD", "Bangladesh"), ("BB", "BRB", "Barbados"),
    ("BY", "BLR", "Belarus"), ("BE", "BEL", "Belgium"), ("BZ", "BLZ", "Belize"),
    ("BJ", "BEN", "Benin"), ("BM", "BMU", "Bermuda"), ("BT", "BTN", "Bhutan"),
    ("BO", "BOL", "Bolivia"), ("BA", "BIH", "Bosnia and Herzegovina"), ("BW", "BWA", "Botswana"),
    ("BR", "BRA", "Brazil"), ("BN", "BRN", "Brunei"), ("BG", "BGR", "Bulgaria"),
    ("BF", "BFA", "Burkina Faso"), ("BI", "BDI", "Burundi"), ("KH", "KHM", "Cambodia"),
    ("CM", "CMR", "Cameroon"), ("CA", "CAN", "Canada"), ("CV", "CPV", "Cape Verde"),
    ("KY", "CYM", "Cayman Islands"), ("CF", "CAF", "Central African Republic"), ("TD", "TCD", "Chad"),
    ("CL", "CHL", "Chile"), ("CN", "CHN", "China"), ("CO", "COL", "Colombia"),
    ("KM", "COM", "Comoros"), ("CG", "COG", "Congo"), ("CD", "COD", "Congo, Democratic Republic of the"),
    ("CR", "CRI", "Costa Rica"), ("CI", "CIV", "Cote d'Ivoire"), ("HR", "HRV", "Croatia"),
    ("CU", "CUB", "Cuba"), ("CW", "CUW", "Curacao"), ("CY", "CYP", "Cyprus"),
    ("CZ", "CZE", "Czech Republic"), ("DK", "DNK", "Denmark"), ("DJ", "DJI", "Djibouti"),
    ("DM", "DMA", "Dominica"), ("DO", "DOM", "Dominican Republic"), ("EC", "ECU", "Ecuador"),
    ("EG", "EGY", "Egypt"), ("SV", "SLV", "El Salvador"), ("GQ", "GNQ", "Equatorial Guinea"),
    ("ER", "ERI", "Eritrea"), ("EE", "EST", "Estonia"), ("SZ", "SWZ", "Eswatini"),
    ("ET", "ETH", "Ethiopia"), ("FJ", "FJI", "Fiji"), ("FI", "FIN", "Finland"),
    ("FR", "FRA", "France"), ("GA", "GAB", "Gabon"), ("GM", "GMB", "Gambia"),
    ("GE", "GEO", "Georgia"), ("DE", "DEU", "Germany"), ("GH", "GHA", "Ghana"),
    ("GI", "GIB", "Gibraltar"), ("GR", "GRC", "Greece"), ("GL", "GRL", "Greenland"),
    ("GD", "GRD", "Grenada"), ("GU", "GUM", "Guam"), ("GT", "GTM", "Guatemala"),
    ("GG", "GGY", "Guernsey"), ("GN", "GIN", "Guinea"), ("GW", "GNB", "Guinea-Bissau"),
    ("GY", "GUY", "Guyana"), ("HT", "HTI", "Haiti"), ("HN", "HND", "Honduras"),
    ("HK", "HKG", "Hong Kong"), ("HU", "HUN", "Hungary"), ("IS", "ISL", "Iceland"),
    ("IN", "IND", "India"), ("ID", "IDN", "Indonesia"), ("IR", "IRN", "Iran"),
    ("IQ", "IRQ", "Iraq"), ("IE", "IRL", "Ireland"), ("IM", "IMN", "Isle of Man"),
    ("IL", "ISR", "Israel"), ("IT", "ITA", "Italy"), ("JM", "JAM", "Jamaica"),
    ("JP", "JPN", "Japan"), ("JE", "JEY", "Jersey"), ("JO", "JOR", "Jordan"),
    ("KZ", "KAZ", "Kazakhstan"), ("KE", "KEN", "Kenya"), ("KI", "KIR", "Kiribati"),
    ("KP", "PRK", "North Korea"), ("KR", "KOR", "South Korea"), ("KW", "KWT", "Kuwait"),
    ("KG", "KGZ", "Kyrgyzstan"), ("LA", "LAO", "Laos"), ("LV", "LVA", "Latvia"),
    ("LB", "LBN", "Lebanon"), ("LS", "LSO", "Lesotho"), ("LR", "LBR", "Liberia"),
    ("LY", "LBY", "Libya"), ("LI", "LIE", "Liechtenstein"), ("LT", "LTU", "Lithuania"),
    ("LU", "LUX", "Luxembourg"), ("MO", "MAC", "Macao"), ("MG", "MDG", "Madagascar"),
    ("MW", "MWI", "Malawi"), ("MY", "MYS", "Malaysia"), ("MV", "MDV", "Maldives"),
    ("ML", "MLI", "Mali"), ("MT", "MLT", "Malta"), ("MH", "MHL", "Marshall Islands"),
    ("MR", "MRT", "Mauritania"), ("MU", "MUS", "Mauritius"), ("MX", "MEX", "Mexico"),
    ("FM", "FSM", "Micronesia"), ("MD", "MDA", "Moldova"), ("MC", "MCO", "Monaco"),
    ("MN", "MNG", "Mongolia"), ("ME", "MNE", "Montenegro"), ("MA", "MAR", "Morocco"),
    ("MZ", "MOZ", "Mozambique"), ("MM", "MMR", "Myanmar"), ("NA", "NAM", "Namibia"),
    ("NR", "NRU", "Nauru"), ("NP", "NPL", "Nepal"), ("NL", "NLD", "Netherlands"),
    ("NC", "NCL", "New Caledonia"), ("NZ", "NZL", "New Zealand"), ("NI", "NIC", "Nicaragua"),
    ("NE", "NER", "Niger"), ("NG", "NGA", "Nigeria"), ("MK", "MKD", "North Macedonia"),
    ("NO", "NOR", "Norway"), ("OM", "OMN", "Oman"), ("PK", "PAK", "Pakistan"),
    ("PW", "PLW", "Palau"), ("PS", "PSE", "Palestine"), ("PA", "PAN", "Panama"),
    ("PG", "PNG", "Papua New Guinea"), ("PY", "PRY", "Paraguay"), ("PE", "PER", "Peru"),
    ("PH", "PHL", "Philippines"), ("PL", "POL", "Poland"), ("PT", "PRT", "Portugal"),
    ("PR", "PRI", "Puerto Rico"), ("QA", "QAT", "Qatar"), ("RO", "ROU", "Romania"),
    ("RU", "RUS", "Russia"), ("RW", "RWA", "Rwanda"), ("KN", "KNA", "Saint Kitts and Nevis"),
    ("LC", "LCA", "Saint Lucia"), ("VC", "VCT", "Saint Vincent and the Grenadines"),
    ("WS", "WSM", "Samoa"), ("SM", "SMR", "San Marino"), ("ST", "STP", "Sao Tome and Principe"),
    ("SA", "SAU", "Saudi Arabia"), ("SN", "SEN", "Senegal"), ("RS", "SRB", "Serbia"),
    ("SC", "SYC", "Seychelles"), ("SL", "SLE", "Sierra Leone"), ("SG", "SGP", "Singapore"),
    ("SK", "SVK", "Slovakia"), ("SI", "SVN", "Slovenia"), ("SB", "SLB", "Solomon Islands"),
    ("SO", "SOM", "Somalia"), ("ZA", "ZAF", "South Africa"), ("SS", "SSD", "South Sudan"),
    ("ES", "ESP", "Spain"), ("LK", "LKA", "Sri Lanka"), ("SD", "SDN", "Sudan"),
    ("SR", "SUR", "Suriname"), ("SE", "SWE", "Sweden"), ("CH", "CHE", "Switzerland"),
    ("SY", "SYR", "Syria"), ("TW", "TWN", "Taiwan"), ("TJ", "TJK", "Tajikistan"),
    ("TZ", "TZA", "Tanzania"), ("TH", "THA", "Thailand"), ("TL", "TLS", "Timor-Leste"),
    ("TG", "TGO", "Togo"), ("TO", "TON", "Tonga"), ("TT", "TTO", "Trinidad and Tobago"),
    ("TN", "TUN", "Tunisia"), ("TR", "TUR", "Turkey"), ("TM", "TKM", "Turkmenistan"),
    ("TV", "TUV", "Tuvalu"), ("UG", "UGA", "Uganda"), ("UA", "UKR", "Ukraine"),
    ("AE", "ARE", "United Arab Emirates"), ("GB", "GBR", "United Kingdom"), ("US", "USA", "United States"),
    ("UY", "URY", "Uruguay"), ("UZ", "UZB", "Uzbekistan"), ("VU", "VUT", "Vanuatu"),
    ("VA", "VAT", "Vatican City"), ("VE", "VEN", "Venezuela"), ("VN", "VNM", "Vietnam"),
    ("VG", "VGB", "British Virgin Islands"), ("VI", "VIR", "U.S. Virgin Islands"), ("YE", "YEM", "Yemen"),
    ("ZM", "ZMB", "Zambia"), ("ZW", "ZWE", "Zimbabwe"),
]

# Colloquial / historical name variants not covered by the primary ISO name above.
_EXTRA_COUNTRY_ALIASES: Dict[str, str] = {
    "uk": "GBR", "great britain": "GBR", "britain": "GBR",
    "united states of america": "USA", "america": "USA",
    "prc": "CHN", "mainland china": "CHN",
    "korea": "KOR", "republic of korea": "KOR", "dprk": "PRK",
    "the netherlands": "NLD", "holland": "NLD",
    "ivory coast": "CIV", "czechia": "CZE", "swaziland": "SWZ", "burma": "MMR",
    "dr congo": "COD", "drc": "COD", "zaire": "COD", "republic of the congo": "COG",
    "macedonia": "MKD", "fyrom": "MKD", "vatican": "VAT", "holy see": "VAT",
    "russian federation": "RUS", "syrian arab republic": "SYR", "lao pdr": "LAO",
    "brunei darussalam": "BRN", "viet nam": "VNM", "east timor": "TLS",
    "palestinian territories": "PSE",
    # Not part of the official ISO 3166-1 list (disputed status), but
    # SAP SuccessFactors ships it as a selectable country/CSF value in
    # real EC config workbooks — verified against production data.
    "kosovo": "XKX", "xkx": "XKX",
}


def _build_country_aliases() -> Dict[str, str]:
    aliases: Dict[str, str] = {}
    for alpha2, alpha3, name in _ISO_COUNTRIES:
        aliases[name.lower()] = alpha3
        aliases[alpha2.lower()] = alpha3
        aliases[alpha3.lower()] = alpha3
    aliases.update(_EXTRA_COUNTRY_ALIASES)
    return aliases


# Mapping: country name variants → ISO alpha-3 / alpha-2
COUNTRY_ALIASES: Dict[str, str] = _build_country_aliases()

# All ISO alpha-3 codes as a single alternation, for spotting an embedded
# country code in a sheet tab name (e.g. "DEU_Payroll") — shared by
# CSF_TAB_RE below and the tab-name fallback in detect_countries_in_workbooks.
_ISO3_ALTERNATION = '|'.join(sorted({alpha3 for _, alpha3, _ in _ISO_COUNTRIES}))
ISO3_TAB_RE = re.compile(r'\b(' + _ISO3_ALTERNATION + r')\b')

# Tab name patterns that definitively indicate CSF (country-specific) sheets
CSF_TAB_PATTERNS = [
    r'\bCSF\b', r'_CSF', r'CSF_',
    r'\bLegal[\s_-]?Entity\b', r'\bPay[\s_-]?Info\b',
    r'\bPayroll\b', r'\bTax[\s_-]?',
    r'\bBenefits?\b', r'\bTime[\s_-]?Off\b',
    r'\bLeave[\s_-]?',
    # ISO 3-letter codes in tab name
    r'\b(?:' + _ISO3_ALTERNATION + r')\b',
]
CSF_TAB_RE = re.compile('|'.join(CSF_TAB_PATTERNS), re.IGNORECASE)

# Column header names that indicate a country filter column
COUNTRY_COL_HEADERS = {
    "country", "country code", "country group", "country/region", "countryofcompany",
    "country of company", "country_code", "country_group", "csf country",
    "legal entity country", "payroll country", "geo", "region", "locale",
    "country (iso)", "iso country", "country iso", "cc",
}

# Tab names that are always GLOBAL (never CSF-filtered)
GLOBAL_TAB_KEYWORDS = {
    "foundation", "corporate", "job_code", "job code", "job classification",
    "pay_grade", "pay grade", "pay_group", "pay group", "pay component",
    "org_chart", "org chart", "position", "department", "division", "cost_center",
    "cost center", "business_unit", "business unit", "location", "workflow",
    "dynamic_group", "dynamic group", "event_reason", "event reason",
    "employment_type", "employment type", "holiday_calendar", "holiday calendar",
    "work_schedule", "work schedule", "global", "template", "instructions",
    "readme", "change_log", "change log", "version", "legend", "overview",
    "table of contents", "toc", "index",
}


# ── File-type text extractors ────────────────────────────────────────────────

def _extract_pptx_text(raw_bytes: bytes) -> str:
    """Extract all visible text from a .pptx file (slides + notes + tables)."""
    try:
        from pptx import Presentation  # python-pptx
        from pptx.util import Inches
        prs = Presentation(io.BytesIO(raw_bytes))
        lines = []
        for slide_num, slide in enumerate(prs.slides, 1):
            lines.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                # Text frames (titles, body, text boxes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            lines.append("\t".join(cells))
                # Notes (often contain country scope in project decks)
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            lines.append(f"[Notes] {notes_text}")
                except Exception:
                    pass
        return "\n".join(lines)[:20000]
    except ImportError:
        return _extract_unicode_strings(raw_bytes)
    except Exception as e:
        print(f"[workbook_processor] pptx extract error: {e}")
        return _extract_unicode_strings(raw_bytes)


def _extract_ppt_text(raw_bytes: bytes) -> str:
    """
    Extract text from an old binary .ppt file.
    python-pptx does not support .ppt; we pull UTF-16 LE and ASCII strings
    directly from the binary compound document, which reliably captures
    slide text, notes, and table cells stored by PowerPoint 97-2003.
    """
    return _extract_unicode_strings(raw_bytes)


def _extract_mpp_text(raw_bytes: bytes) -> str:
    """
    Extract human-readable text from a Microsoft Project .mpp file.

    .mpp is a binary OLE Compound Document. Task names, resource names,
    notes, and custom fields are stored as UTF-16 LE strings. Countries
    appear in task names ("Germany - Deploy"), milestones, or resource pools.

    Strategy:
      1. Try python-mpxj (Java-based wrapper) for structured extraction.
      2. Fall back to Unicode/ASCII string mining if Java/mpxj not available.
    """
    # Attempt structured extraction via mpxj if available
    try:
        import mpxj  # pip install mpxj (requires Java)
        import tempfile, os as _os
        with tempfile.NamedTemporaryFile(suffix=".mpp", delete=False) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name
        try:
            from mpxj import ProjectReader
            project = ProjectReader().read(tmp_path)
            lines = []
            for task in project.tasks:
                if task.name:
                    lines.append(task.name)
                if task.notes:
                    lines.append(task.notes)
            for res in project.resources:
                if res.name:
                    lines.append(res.name)
            return "\n".join(lines)[:20000]
        finally:
            _os.unlink(tmp_path)
    except Exception:
        pass

    # Fallback: mine raw Unicode strings — works well for task/milestone names
    return _extract_unicode_strings(raw_bytes)


def _extract_docx_text(raw_bytes: bytes) -> str:
    """Extract paragraph text from a .docx file via its internal XML."""
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        lines = []
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as zf:
            for name in zf.namelist():
                if name.startswith("word/") and name.endswith(".xml"):
                    xml_bytes = zf.read(name)
                    root = ET.fromstring(xml_bytes)
                    for para in root.iter(f"{{{NS}}}p"):
                        texts = [t.text or "" for t in para.iter(f"{{{NS}}}t")]
                        line = "".join(texts).strip()
                        if line:
                            lines.append(line)
        return "\n".join(lines)[:20000]
    except Exception as e:
        print(f"[workbook_processor] docx extract error: {e}")
        return raw_bytes.decode("utf-8", errors="replace")[:15000]


def _extract_unicode_strings(raw_bytes: bytes, min_len: int = 4) -> str:
    """
    Mine printable ASCII and UTF-16 LE strings from arbitrary binary data.
    Effective for .ppt, .mpp, and other OLE compound document formats where
    country names, task names, and region labels appear as inline text.
    """
    found = set()

    # ASCII printable runs
    for m in re.finditer(rb'[ -~]{' + str(min_len).encode() + rb',}', raw_bytes):
        s = m.group().decode("ascii", errors="ignore").strip()
        if s:
            found.add(s)

    # UTF-16 LE runs (Windows-native string encoding in OLE/binary Office files)
    for m in re.finditer(rb'(?:[\x20-\x7e]\x00){' + str(min_len).encode() + rb',}', raw_bytes):
        try:
            s = m.group().decode("utf-16-le", errors="ignore").strip()
            if s:
                found.add(s)
        except Exception:
            pass

    # Sort longer strings first (task names tend to be longer than noise)
    ordered = sorted(found, key=len, reverse=True)
    return "\n".join(ordered)[:20000]


# Canonical display name per ISO-3 code (for countries detected directly from
# workbook data, where we only have a code/alias and need something to show).
# Derived from the same master list COUNTRY_ALIASES is built from, so every
# code that can be *detected* also has a display name.
ISO3_TO_NAME: Dict[str, str] = {alpha3: name for _, alpha3, name in _ISO_COUNTRIES}
ISO3_TO_NAME["XKX"] = "Kosovo"  # see _EXTRA_COUNTRY_ALIASES — not in ISO 3166-1


# ── Helpers ───────────────────────────────────────────────────────────────────

# Closing paren is optional: a real production "CSF Job Info" sheet has
# "Tunisia (TUN" (no closing paren) for every Tunisia row, while other
# sheets in the same workbook have the well-formed "Tunisia (TUN)" —
# an inconsistency in the source data itself. Requiring the closing
# paren meant every Tunisia row in that one sheet silently failed to
# resolve and was kept regardless of country selection.
_NAME_CODE_RE = re.compile(r'^(.*?)\s*\(([A-Za-z]{2,3})\)?\s*$')


def _normalise_country(raw: str) -> Optional[str]:
    """
    Return ISO-3 code for a country string, or None if unrecognised.

    Real SF EC workbooks' CSF country columns commonly format values as
    "Angola (AGO)", "Argentina (ARG)" — name plus code in parentheses,
    not a bare name or code (verified directly against the production
    "National ID" sheet). Without this, exact-match lookup fails for
    every such value, silently excluding all CSF data regardless of
    which countries are actually in scope.
    """
    key = raw.strip().lower()
    direct = COUNTRY_ALIASES.get(key)
    if direct:
        return direct

    m = _NAME_CODE_RE.match(raw.strip())
    if m:
        name_part, code_part = m.group(1).strip().lower(), m.group(2).strip().lower()
        return COUNTRY_ALIASES.get(code_part) or COUNTRY_ALIASES.get(name_part)
    return None


def _resolve_iso3(value: Any) -> Optional[str]:
    """Resolve a raw cell value (name, alias, or ISO code) to ISO-3, or None."""
    if value is None:
        return None
    s = str(value).strip()
    if not s:
        return None
    return _normalise_country(s)


def _is_csf_tab_by_name(sheet_name: str) -> bool:
    """Heuristic: does this tab name look like a CSF (country-specific) sheet?"""
    low = sheet_name.lower()
    for kw in GLOBAL_TAB_KEYWORDS:
        if kw in low:
            return False
    return bool(CSF_TAB_RE.search(sheet_name))


def _find_country_column(ws) -> Optional[int]:
    """
    Scan the first 5 rows of a sheet to find the column index (1-based)
    whose header most likely represents a country/country-code field.
    Returns None if not found.

    Uses bounded iter_rows() rather than ws.cell(row=, column=) random
    access — the latter isn't supported on read-only worksheets, which
    is how sheets are loaded (see load_workbook_from_path): eager-mode
    loading of a real SF EC workbook was measured at 250-400s, versus
    15-30s in read-only mode.
    """
    for row in ws.iter_rows(min_row=1, max_row=5):
        for cell in row:
            if cell.value and isinstance(cell.value, str):
                header = cell.value.strip().lower()
                if header in COUNTRY_COL_HEADERS:
                    return cell.column
    return None


def _get_header_row(ws) -> int:
    """
    Return the 1-based row index of the header row (usually 1 or 2).

    Row index comes from enumerate(), not any individual cell's .row —
    a row's first element can be a lightweight EmptyCell placeholder
    (read-only mode's stand-in for a blank position within a row's
    populated range), which doesn't carry a .row attribute.
    """
    for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=4), start=1):
        non_empty = [c.value for c in row[:10] if c.value is not None and str(c.value).strip()]
        if len(non_empty) >= 2:
            return row_idx
    return 1


def _stream_real_rows(ws, max_gap: int = 10000):
    """
    Yield (row_idx, [(col_idx, cell), ...]) only for rows that actually
    contain at least one populated (value or styled) cell, reading the
    sheet forward and stopping early once `max_gap` consecutive fully-
    empty rows have been seen.

    Real-world Excel workbooks routinely have a handful of cells stranded
    far outside their real content block — e.g. one leftover cell near
    the sheet's absolute row limit, which genuinely exists but drags
    ws.max_row out past a million, even though real data ends a few
    dozen or few thousand rows in. Naively iterating to ws.max_row would
    mean streaming through — or, in eager/non-read-only mode, having
    openpyxl backfill a dense grid of Cell objects for — over a million
    phantom rows. The early-break-on-gap approach avoids ever reaching
    that phantom range at all, for both read-only (cheap either way,
    just avoids wasted iterations) and eager worksheets (avoids the
    expensive backfill entirely).

    Verified against a real SF EC workbook: a legitimately dense
    40,008-row sheet has zero gaps over 50 rows (so is read in full),
    while a sheet with real content ending at row ~1,062 had a single
    stray cell at row 1,048,573 — a >1,000,000-row gap that's cut here.

    Row index is read from a populated cell's own .row, not derived
    from enumerate() position. ws.iter_rows() (no min_row given) always
    yields starting from actual row 1 — including leading fully-empty
    rows — regardless of where ws.min_row says real content begins;
    labeling yielded rows by counting from ws.min_row (as an earlier
    version of this function did) silently mislabels every row by that
    offset whenever a sheet's real content doesn't start at row 1
    (confirmed: correct kept/removed *counts* but the wrong physical
    rows deleted). Read-only mode's lightweight EmptyCell placeholders
    for blank positions within a row lack .row entirely, which is why
    this doesn't just take row[0].row — but every cell that lands in
    `populated` below is guaranteed to be a real Cell/ReadOnlyCell
    (EmptyCell always has value=None and no has_style attribute, so it
    can never satisfy the filter condition), so populated[0][1].row is
    always safe.
    """
    empty_streak = 0
    for row in ws.iter_rows():
        populated = [(c.column, c) for c in row if c.value is not None or getattr(c, "has_style", False)]
        if populated:
            empty_streak = 0
            yield populated[0][1].row, populated
        else:
            empty_streak += 1
            if empty_streak > max_gap:
                return


# ── Main public functions ────────────────────────────────────────────────────

async def extract_countries_from_document(file_data_b64: str, file_name: str) -> List[Dict[str, str]]:
    """
    Use Claude to extract all in-scope countries from a SOW or SP51 document.

    Returns a list of dicts: [{"name": "Germany", "iso3": "DEU"}, ...]
    """
    try:
        raw_bytes = base64.b64decode(file_data_b64)
    except Exception:
        return []

    ext = Path(file_name).suffix.lower()

    # Build the document content block(s) for Claude based on file type
    doc_content_blocks = []

    if ext in ('.xlsx', '.xls'):
        try:
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames[:5]:
                ws = wb[sheet_name]
                sheet_text = f"\n[Sheet: {sheet_name}]\n"
                for row in ws.iter_rows(max_row=200, values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        sheet_text += row_str + "\n"
                text_parts.append(sheet_text)
            combined = "\n".join(text_parts)[:20000]
            doc_content_blocks = [{"type": "text", "text": combined}]
        except Exception as e:
            doc_content_blocks = [{"type": "text", "text": f"[Could not parse Excel: {e}]"}]

    elif ext == '.pdf':
        # Claude natively reads PDFs — send as document block for best accuracy
        doc_content_blocks = [
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": file_data_b64,
                },
            }
        ]

    elif ext == '.docx':
        text = _extract_docx_text(raw_bytes)
        doc_content_blocks = [{"type": "text", "text": text}]

    elif ext == '.pptx':
        text = _extract_pptx_text(raw_bytes)
        doc_content_blocks = [{"type": "text", "text": f"[PowerPoint Presentation: {file_name}]\n\n{text}"}]

    elif ext == '.ppt':
        text = _extract_ppt_text(raw_bytes)
        doc_content_blocks = [{"type": "text", "text": f"[PowerPoint 97-2003: {file_name}]\n\n{text}"}]

    elif ext == '.mpp':
        text = _extract_mpp_text(raw_bytes)
        doc_content_blocks = [{"type": "text", "text": f"[Microsoft Project File: {file_name}]\n\n{text}"}]

    else:
        # .txt, .csv, and any other text-based format
        try:
            text = raw_bytes.decode("utf-8", errors="replace")[:20000]
        except Exception:
            text = f"[Binary file: {file_name}]"
        doc_content_blocks = [{"type": "text", "text": text}]

    client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))

    system = (
        "You are an expert SAP SuccessFactors implementation consultant. "
        "You are reading a Statement of Work (SOW) or SP51 scope document. "
        "Your task is to extract ALL countries that are explicitly in scope for this implementation. "
        "Return ONLY a JSON array of objects with 'name' (full English country name) and 'iso3' (ISO 3166-1 alpha-3 code). "
        "Example: [{\"name\": \"Germany\", \"iso3\": \"DEU\"}, {\"name\": \"United Kingdom\", \"iso3\": \"GBR\"}]. "
        "Do not include countries that are mentioned as out-of-scope, excluded, or future phases. "
        "If the document mentions regions (e.g., 'EMEA', 'APAC'), list only the specific countries explicitly named. "
        "Return ONLY the JSON array, no other text."
    )

    content_blocks = doc_content_blocks + [
        {
            "type": "text",
            "text": (
                f"Extract all in-scope countries from this {'SOW' if 'sow' in file_name.lower() else 'SP51/scope'} document. "
                "Return a JSON array with 'name' and 'iso3' fields only."
            ),
        }
    ]

    try:
        response = client.messages.create(
            model="claude-sonnet-5",
            max_tokens=1024,
            system=system,
            messages=[{"role": "user", "content": content_blocks}],
        )
        raw_text = response.content[0].text.strip()
        # Strip markdown code fences if present
        raw_text = re.sub(r"^```[a-z]*\n?", "", raw_text)
        raw_text = re.sub(r"\n?```$", "", raw_text)
        countries = json.loads(raw_text)
        if isinstance(countries, list):
            # Validate and enrich
            result = []
            for c in countries:
                if isinstance(c, dict) and c.get("name"):
                    iso3 = c.get("iso3", "")
                    if not iso3:
                        iso3 = _normalise_country(c["name"]) or ""
                    result.append({"name": c["name"].strip(), "iso3": iso3.upper()})
            return result
    except Exception as e:
        print(f"[workbook_processor] country extraction error: {e}")

    return []


def load_workbook_from_path(ref_file_path: str) -> Optional[openpyxl.Workbook]:
    """
    Load an openpyxl workbook from disk. Returns None on failure.

    read_only=True: eager (writable) loading was measured at 250-400s
    against a real ~3.3MB, 32-sheet SF EC workbook, versus 15-30s in
    read-only mode — eager mode builds a full mutable Cell object graph
    with resolved styles for every touched position, which is far more
    expensive than read-only's streaming SAX-based access. Read-only
    cells still expose .value, .font, .fill, .border, .alignment, and
    .has_style (verified directly against this workbook), so nothing
    downstream that only reads cells is affected — only random access
    via ws.cell(row=, column=) is unavailable, which is why
    _find_country_column/_get_header_row/_stream_real_rows all use
    iter_rows() instead.
    """
    try:
        return openpyxl.load_workbook(ref_file_path, data_only=True, read_only=True)
    except Exception as e:
        print(f"[workbook_processor] load error {ref_file_path}: {e}")
        return None


def detect_countries_in_workbooks(slot_files: List[Dict], refs_dir: Path) -> List[Dict[str, Any]]:
    """
    Scan all grounded reference workbooks and return the distinct countries
    actually present in their CSF sheets — for the manual country-selection
    flow, where no SOW is uploaded and countries must be read from the data
    that is already grounded.

    Detection order per CSF sheet:
      1. If a country column was found, resolve every data-row value in it.
      2. Otherwise, the country is usually encoded in the tab name itself
         (e.g. "DEU_Payroll") — check for an embedded ISO-3 code there.
      3. As a last resort, scan the first few cells of each data row.

    Returns: [{"name": "Germany", "iso3": "DEU", "count": 42}, ...]
             sorted alphabetically by name.
    """
    found: Dict[str, int] = {}

    for slot_info in slot_files:
        ref_id = slot_info.get("ref_id", "")
        file_path = None
        for candidate in refs_dir.glob(f"{ref_id}.*"):
            file_path = candidate
            break
        if not file_path or not file_path.exists():
            continue

        wb = load_workbook_from_path(str(file_path))
        if not wb:
            continue

        classification = classify_sheets(wb)
        for sheet_name, info in classification.items():
            if info["type"] != "csf":
                continue
            ws = wb[sheet_name]
            country_col = info.get("country_col")
            header_row  = info.get("header_row", 1)

            if country_col:
                for row_idx, col_cells in _stream_real_rows(ws):
                    if row_idx <= header_row:
                        continue
                    for col_idx, cell in col_cells:
                        if col_idx == country_col:
                            iso3 = _resolve_iso3(cell.value)
                            if iso3:
                                found[iso3] = found.get(iso3, 0) + 1
                            break
            else:
                tab_match = ISO3_TAB_RE.search(sheet_name.upper())
                if tab_match:
                    iso3 = tab_match.group(1)
                    found[iso3] = found.get(iso3, 0) + 1
                else:
                    for row_idx, col_cells in _stream_real_rows(ws):
                        if row_idx <= header_row:
                            continue
                        for col_idx, cell in col_cells:
                            if col_idx > 5:
                                break
                            iso3 = _resolve_iso3(cell.value)
                            if iso3:
                                found[iso3] = found.get(iso3, 0) + 1

    result = [
        {"name": ISO3_TO_NAME.get(iso3, iso3), "iso3": iso3, "count": count}
        for iso3, count in found.items()
    ]
    result.sort(key=lambda c: c["name"])
    return result


def classify_sheets(wb: openpyxl.Workbook) -> Dict[str, Dict]:
    """
    Classify all sheets in a workbook as:
      - 'csf': country-specific, must be filtered
      - 'global': applies to all countries, copy as-is
      - 'meta': instructions / legend / TOC, copy as-is

    Returns dict: sheet_name -> {type, country_col (if csf), header_row}
    """
    result = {}
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        low = sheet_name.lower().strip()

        # Check for meta sheets first
        if any(kw in low for kw in ("readme", "instruction", "legend", "toc", "table of content", "version", "change log", "changelog")):
            result[sheet_name] = {"type": "meta"}
            continue

        # Try to find a country column in the sheet data BEFORE falling back
        # to a keyword match on the tab name. A real detected country column
        # is a far stronger, more specific signal than a substring match —
        # verified against a real SF EC workbook: "CSF Personal Info (Global
        # Info)" contains "global" (from its own "Global Info" suffix) and
        # was being short-circuited into the global bucket before this check
        # ever ran, even though it has the same country-column structure as
        # every other CSF sheet in the file and its name literally starts
        # with "CSF".
        country_col = _find_country_column(ws)
        header_row = _get_header_row(ws)

        if country_col:
            result[sheet_name] = {"type": "csf", "country_col": country_col, "header_row": header_row}
            continue

        # Check if globally excluded by name
        if any(kw in low for kw in GLOBAL_TAB_KEYWORDS):
            result[sheet_name] = {"type": "global"}
            continue

        if _is_csf_tab_by_name(sheet_name):
            # Tab name pattern suggests CSF — scan first data column for country codes
            result[sheet_name] = {"type": "csf", "country_col": None, "header_row": header_row}
        else:
            result[sheet_name] = {"type": "global"}

    return result


def _delete_row_ranges(ws, rows_to_delete: List[int]) -> None:
    """
    Merge consecutive row numbers into contiguous (start, count) ranges
    and delete each range with one ws.delete_rows() call, from the
    bottom up. openpyxl's delete_rows shifts every remaining cell below
    the deletion point on EVERY call — calling it once per individual
    row is O(n^2) for a sheet with many rows to remove (confirmed: this
    made a real ~1,500-row sheet with ~1,460 rows to delete hang for
    10+ minutes). Merging into ranges cuts the call count down to the
    number of contiguous deleted blocks, typically small since related
    rows are themselves usually contiguous in these templates.
    """
    ranges: List[Tuple[int, int]] = []
    for row_idx in rows_to_delete:
        if ranges and row_idx == ranges[-1][0] + ranges[-1][1]:
            start, count = ranges[-1]
            ranges[-1] = (start, count + 1)
        else:
            ranges.append((row_idx, 1))

    for start, count in reversed(ranges):
        ws.delete_rows(start, count)


def _filter_sheet_rows_in_place(ws, info: Dict, iso3_set: set) -> Tuple[int, int]:
    """
    Delete rows from a CSF worksheet that belong to an out-of-scope
    country, in place — nothing else about the sheet (styles, merged
    cells, data validations, conditional formatting, frozen panes, row
    heights, column widths) is ever touched, so it stays exactly as
    authored in the grounded file.

    A row is only ever a delete candidate if its country-column cell
    actually resolves to a recognised country — if it doesn't (blank,
    a column header/label, section documentation, anything else that
    isn't naming a specific country), the row is always kept untouched.
    This is deliberately NOT gated on header_row: real SF EC CSF sheets
    can have a more complex layout than "header row(s), then immediate
    data" — e.g. a field-definition/documentation block before the
    actual per-country data table starts much further down. Using a
    single guessed header_row cutoff as "everything after this is
    filterable data" was deleting that documentation block wholesale
    (verified against this exact shape locally). Driving the decision
    off whether the cell resolves to a country at all sidesteps needing
    to know where any such boundary actually falls.

    Returns (kept_count, removed_count).
    """
    country_col = info.get("country_col")

    rows_to_delete: List[int] = []
    kept = 0

    for row_idx, col_cells in _stream_real_rows(ws):
        resolved: Optional[str] = None
        if country_col:
            for col_idx, cell in col_cells:
                if col_idx == country_col:
                    resolved = _resolve_iso3(cell.value)
                    break
        else:
            # No explicit country column found — check the first 5
            # columns of the row for any value that resolves to a
            # country (conservative: the first hit wins).
            for col_idx, cell in col_cells:
                if col_idx > 5:
                    break
                resolved = _resolve_iso3(cell.value)
                if resolved:
                    break

        if resolved is None:
            # Not a country-specific data row — leave it exactly as-is.
            kept += 1
        elif resolved in iso3_set:
            kept += 1
        else:
            rows_to_delete.append(row_idx)

    _delete_row_ranges(ws, rows_to_delete)
    return kept, len(rows_to_delete)


# Picklist IDs where the row's value column names the country directly
# (rather than the ID itself carrying a country suffix) — see
# _filter_picklists_sheet_in_place.
_COUNTRY_MASTER_PICKLISTS = {"isocountrylist", "csfcountry", "country"}
_PICKLIST_ID_SUFFIX_RE = re.compile(r'^(.+)_([A-Za-z]{3})$')


def _filter_picklists_sheet_in_place(ws, iso3_set: set) -> Tuple[int, int]:
    """
    Filter the "Picklists" reference sheet to the in-scope countries.

    This sheet has a different shape than CSF sheets — there's no
    single country column. Verified directly against a real ~40,000-row
    production "Picklists" sheet:

    - The large majority of rows (verified: 35,760 of ~40,000) encode
      the country as a suffix on the Picklist ID column itself, e.g.
      "ACADEMICDEGREE_ARE" = the Academic Degree picklist's UAE variant
      — not as a separate column value. Confirmed systematic (194
      distinct base picklist names, each with 1-50 country variants),
      not coincidental.
    - A handful of specific picklists — ISOCountryList, csfCountry,
      country — are themselves the master list of selectable countries:
      every row has the SAME Picklist ID, and the country is named in
      the External Code / Description columns instead.
    - Everything else (Nationality, state, jobRegion, currency,
      nameprefix, and single-value global picklists like BLOODGROUP)
      has no reliable per-row country signal — e.g. Nationality rows
      are demonyms ("Estonian") with no clean adjective-to-country
      mapping, and state/jobRegion mix country-qualified and bare
      values inconsistently. These are left untouched: wrongly
      deleting them would remove valid global reference data, a worse
      failure than leaving a few extra picklist rows in place.

    Returns (kept_count, removed_count).
    """
    rows_to_delete: List[int] = []
    kept = 0

    for row_idx, col_cells in _stream_real_rows(ws):
        col_map = {c: cell for c, cell in col_cells}
        picklist_id_cell = col_map.get(3)
        picklist_id = picklist_id_cell.value if picklist_id_cell else None

        resolved: Optional[str] = None
        if isinstance(picklist_id, str) and picklist_id.strip():
            base_lower = picklist_id.strip().lower()
            if base_lower in _COUNTRY_MASTER_PICKLISTS:
                ext_cell  = col_map.get(4)
                desc_cell = col_map.get(5)
                resolved = (
                    (_resolve_iso3(ext_cell.value) if ext_cell else None)
                    or (_resolve_iso3(desc_cell.value) if desc_cell else None)
                )
            else:
                m = _PICKLIST_ID_SUFFIX_RE.match(picklist_id.strip())
                if m:
                    resolved = _resolve_iso3(m.group(2))

        if resolved is None:
            # No reliable country signal — leave it exactly as-is.
            kept += 1
        elif resolved in iso3_set:
            kept += 1
        else:
            rows_to_delete.append(row_idx)

    _delete_row_ranges(ws, rows_to_delete)
    return kept, len(rows_to_delete)


def _add_summary_sheet(out_wb: openpyxl.Workbook, in_scope_countries: List[Dict[str, str]], source_count: int) -> None:
    summary = out_wb.create_sheet("_Summary", 0)
    summary["A1"] = "Configuration Workbook — Country-Filtered Output"
    summary["A1"].font = Font(bold=True, size=14)
    summary["A3"] = "In-Scope Countries:"
    summary["A3"].font = Font(bold=True)
    for i, c in enumerate(in_scope_countries, start=4):
        summary[f"A{i}"] = c.get("name", "")
        summary[f"B{i}"] = c.get("iso3", "")
    summary.column_dimensions["A"].width = 30
    summary.column_dimensions["B"].width = 10

    import datetime as dt_mod
    summary[f"A{len(in_scope_countries)+5}"] = f"Generated: {dt_mod.datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
    summary[f"A{len(in_scope_countries)+6}"] = f"Source workbooks: {source_count}"


def _find_slot_file(slot_info: Dict, refs_dir: Path) -> Optional[Path]:
    ref_id = slot_info.get("ref_id", "")
    for candidate in refs_dir.glob(f"{ref_id}.*"):
        return candidate
    return None


_VERSION_MARKER_RE = re.compile(r'\bV\d+\b', re.IGNORECASE)


def _clean_output_basename(stem: str) -> str:
    """
    Normalize a grounded file's name for use in the generated output
    filename: strips a "myConcerto_" template prefix and replaces any
    version marker (V1, V2, V10, ...) with "v1" — the output is always
    a fresh v1 artifact of that generation, independent of whatever
    version the source template itself happened to be.
    """
    name = re.sub(r'^myConcerto_', '', stem, flags=re.IGNORECASE)
    name = _VERSION_MARKER_RE.sub('v1', name)
    return name


def load_workbook_for_editing(ref_file_path: str) -> Optional[openpyxl.Workbook]:
    """
    Load a workbook in fully writable (eager) mode — required for the
    generate/filter path, which deletes disqualified CSF rows in place
    on the loaded object rather than rebuilding a new workbook cell by
    cell. That guarantees exact fidelity (styles, merged cells, data
    validations, conditional formatting, frozen panes, tab colors, row
    heights, formulas) since nothing but the removed rows is ever
    touched — this file is never edited on disk either; the original
    ref file stays untouched and only this in-memory copy is modified,
    then saved to a new output buffer.

    data_only=False preserves formulas rather than collapsing them to
    their last cached value. The detection-only loader
    (load_workbook_from_path) uses data_only=True + read_only=True on
    purpose, since it only ever reads values and never saves the
    result — this loader is for the path that does.

    Trade-off: eager loading was measured at 250-400s against a real
    ~3.3MB, 32-sheet SF EC workbook, versus 15-30s in read-only mode.
    That cost is accepted here in exchange for fidelity, which is why
    generation runs as a background job with progress streaming
    (run_workbook_generation) instead of a single blocking request.
    """
    try:
        return openpyxl.load_workbook(ref_file_path, data_only=False, read_only=False)
    except Exception as e:
        print(f"[workbook_processor] eager load error {ref_file_path}: {e}")
        return None


# ═══════════════════════════════════════════════════════════════════════════
# Background job orchestration + SSE progress for workbook generation
# ═══════════════════════════════════════════════════════════════════════════
#
# Mirrors the in-memory SSE queue pattern the cascade feature uses (see
# context_store.py) but is fully independent of it — this module owns
# its own job/queue/result state; nothing here touches cascade code.

_workbook_sse_queues: Dict[str, asyncio.Queue] = {}
_workbook_job_results: "OrderedDict[str, Dict[str, Any]]" = OrderedDict()
_MAX_RETAINED_JOB_RESULTS = 5


def create_workbook_job(job_id: str) -> asyncio.Queue:
    q: asyncio.Queue = asyncio.Queue(maxsize=1000)
    _workbook_sse_queues[job_id] = q
    return q


def get_workbook_queue(job_id: str) -> Optional[asyncio.Queue]:
    return _workbook_sse_queues.get(job_id)


async def emit_workbook_event(job_id: str, event_type: str, data: Dict[str, Any]) -> None:
    q = _workbook_sse_queues.get(job_id)
    if q:
        try:
            await asyncio.wait_for(q.put({"event": event_type, "data": data}), timeout=2.0)
        except (asyncio.TimeoutError, asyncio.QueueFull):
            pass  # Non-blocking — never block generation on a slow/gone client


def end_workbook_stream(job_id: str) -> None:
    q = _workbook_sse_queues.get(job_id)
    if q:
        try:
            q.put_nowait(None)
        except asyncio.QueueFull:
            pass


def cleanup_workbook_queue(job_id: str) -> None:
    """Remove the queue after the SSE stream closes."""
    _workbook_sse_queues.pop(job_id, None)


def store_workbook_result(job_id: str, slots: List[Dict[str, Any]]) -> None:
    """
    Keep generated file bytes in memory, keyed by job_id, for the
    download endpoints. Evicts the oldest job once more than
    _MAX_RETAINED_JOB_RESULTS are held — results aren't persisted to
    disk/DB (consistent with the cascade feature's in-memory-only
    session state; this is a low-traffic admin tool).
    """
    _workbook_job_results[job_id] = {"slots": slots}
    while len(_workbook_job_results) > _MAX_RETAINED_JOB_RESULTS:
        _workbook_job_results.popitem(last=False)


def get_workbook_result(job_id: str) -> Optional[Dict[str, Any]]:
    return _workbook_job_results.get(job_id)


def _generate_one_slot_sync(
    slot_info: Dict,
    in_scope_countries: List[Dict[str, str]],
    refs_dir: Path,
    progress_cb,
) -> Optional[bytes]:
    """
    Synchronous, CPU-bound per-slot pipeline: load eagerly, filter CSF
    sheets in place, save. Runs inside asyncio.to_thread; progress_cb
    is safe to call from this worker thread (see run_workbook_generation).
    """
    slot_num = slot_info.get("slot", 1)

    file_path = _find_slot_file(slot_info, refs_dir)
    if not file_path or not file_path.exists():
        progress_cb("slot_error", {"slot": slot_num, "message": "Reference file not found on disk"})
        return None

    progress_cb("slot_loading", {"slot": slot_num, "file_name": slot_info.get("file_name", "")})
    wb = load_workbook_for_editing(str(file_path))
    if not wb:
        progress_cb("slot_error", {"slot": slot_num, "message": "Could not open workbook — file may be corrupt"})
        return None

    classification = classify_sheets(wb)
    csf_sheet_names = [name for name, info in classification.items() if info["type"] == "csf"]
    progress_cb("slot_classified", {
        "slot": slot_num,
        "total_sheets": len(wb.sheetnames),
        "csf_sheet_count": len(csf_sheet_names),
        "csf_sheets": csf_sheet_names,
    })

    iso3_set = {c.get("iso3", "").upper() for c in in_scope_countries if c.get("iso3")}

    for sheet_name in csf_sheet_names:
        progress_cb("sheet_start", {"slot": slot_num, "sheet": sheet_name})
        ws = wb[sheet_name]
        kept, removed = _filter_sheet_rows_in_place(ws, classification[sheet_name], iso3_set)
        progress_cb("sheet_done", {"slot": slot_num, "sheet": sheet_name, "kept": kept, "removed": removed})
        print(f"  [CSF filter] slot {slot_num} — {sheet_name}: kept {kept}, removed {removed}")

    # The "Picklists" reference sheet has a different shape than CSF
    # sheets (no single country column — see _filter_picklists_sheet_in_place)
    # so it isn't caught by classify_sheets' csf/global split and needs
    # its own pass, matched by name since it's a standard SF EC tab.
    picklists_name = next((n for n in wb.sheetnames if n.strip().lower() == "picklists"), None)
    if picklists_name:
        progress_cb("sheet_start", {"slot": slot_num, "sheet": picklists_name})
        ws = wb[picklists_name]
        kept, removed = _filter_picklists_sheet_in_place(ws, iso3_set)
        progress_cb("sheet_done", {"slot": slot_num, "sheet": picklists_name, "kept": kept, "removed": removed})
        print(f"  [Picklists filter] slot {slot_num} — {picklists_name}: kept {kept}, removed {removed}")

    _add_summary_sheet(wb, in_scope_countries, source_count=1)

    progress_cb("slot_saving", {"slot": slot_num})
    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    progress_cb("slot_done", {
        "slot": slot_num,
        "file_name": slot_info.get("file_name", ""),
        "size_kb": round(len(xlsx_bytes) / 1024, 1),
    })
    return xlsx_bytes


async def run_workbook_generation(
    job_id: str,
    in_scope_countries: List[Dict[str, str]],
    slot_files: List[Dict],
    refs_dir: Path,
) -> None:
    """
    Background orchestrator for a workbook generation job: processes
    each grounded slot in turn — sequentially, since each slot's eager
    load is itself CPU-heavy enough that parallelizing would just
    contend for the same core — streaming progress via SSE, and
    produces one output file per slot (never merged: each grounded
    workbook keeps its own identity and formatting).
    """
    loop = asyncio.get_running_loop()

    def progress_cb(event_type: str, data: Dict[str, Any]) -> None:
        asyncio.run_coroutine_threadsafe(emit_workbook_event(job_id, event_type, data), loop)

    country_labels = ", ".join(c.get("iso3", c.get("name", "")) for c in in_scope_countries)
    await emit_workbook_event(job_id, "job_start", {
        "total_slots": len(slot_files),
        "countries": in_scope_countries,
    })
    print(f"[workbook_processor] Job {job_id}: generating for {country_labels} across {len(slot_files)} slot(s)")

    results: List[Dict[str, Any]] = []
    try:
        for slot_info in slot_files:
            slot_num = slot_info.get("slot", 1)
            await emit_workbook_event(job_id, "slot_start", {
                "slot": slot_num,
                "file_name": slot_info.get("file_name", ""),
            })

            xlsx_bytes = await asyncio.to_thread(
                _generate_one_slot_sync, slot_info, in_scope_countries, refs_dir, progress_cb,
            )

            if xlsx_bytes:
                country_codes = "_".join(c.get("iso3", "") for c in in_scope_countries[:5])
                base_name = _clean_output_basename(Path(slot_info.get("file_name", f"Workbook_{slot_num}")).stem)
                out_name = f"{base_name}_{country_codes}.xlsx"
                results.append({
                    "slot": slot_num,
                    "file_name": out_name,
                    "size_kb": round(len(xlsx_bytes) / 1024, 1),
                    "bytes": xlsx_bytes,
                })

        store_workbook_result(job_id, results)
        await emit_workbook_event(job_id, "job_complete", {
            "slots": [{"slot": r["slot"], "file_name": r["file_name"], "size_kb": r["size_kb"]} for r in results],
            "zip_available": len(results) > 1,
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        await emit_workbook_event(job_id, "job_error", {"message": str(e)})
    finally:
        end_workbook_stream(job_id)
