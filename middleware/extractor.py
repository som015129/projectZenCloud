"""
extractor.py — File text extraction for all supported formats.
Used both for input files (to feed Claude) and generated files (to build content_summary).
FIX A preserved: XML/txt/csv return raw text directly.
"""

import os
from typing import Optional


def extract_text(file_path: str, ext: Optional[str] = None) -> str:
    """
    Extract plain text from a file based on its extension.
    Returns empty string on failure — never raises.
    """
    if not ext:
        ext = os.path.splitext(file_path)[1].lstrip(".").lower()
    ext = ext.lower().strip(".")

    try:
        # ── DOCX ──────────────────────────────────────────────────────────
        if ext in ("docx", "doc"):
            import mammoth
            with open(file_path, "rb") as f:
                result = mammoth.extract_raw_text(f)
            return (result.value or "")[:80000]

        # ── XLSX / XLS ────────────────────────────────────────────────────
        if ext in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(",".join(cells))
                parts.append(f"[{name}]\n" + "\n".join(rows))
            return "\n\n".join(parts)[:80000]

        # ── PPTX ──────────────────────────────────────────────────────────
        if ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(file_path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)[:80000]

        # ── PDF ───────────────────────────────────────────────────────────
        if ext == "pdf":
            text = ""
            # Primary: pdfminer.six
            try:
                from pdfminer.high_level import extract_text as pdfminer_extract
                text = (pdfminer_extract(file_path) or "").strip()
            except Exception as e:
                print(f"   ⚠ pdfminer failed ({e}), trying pypdf")

            # Fallback: pypdf
            if not text:
                try:
                    import pypdf
                    reader = pypdf.PdfReader(file_path)
                    pages = []
                    for page in reader.pages:
                        t = page.extract_text()
                        if t:
                            pages.append(t)
                    text = "\n".join(pages).strip()
                except Exception as e:
                    print(f"   ⚠ pypdf failed ({e})")

            return text[:80000]

        # ── FIX A: XML / TXT / CSV — return raw text ──────────────────────
        if ext in ("xml", "txt", "csv", "md", "json"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()[:80000]

    except Exception as e:
        print(f"   ⚠ extractor error ({ext}): {e}")

    return ""


def extract_from_bytes(content: bytes, ext: str) -> str:
    """
    Extract text from raw bytes (used for uploaded files sent as base64).
    Writes to a temp file, extracts, cleans up.
    """
    import tempfile

    ext = ext.lower().strip(".")
    suffix = f".{ext}"

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        return extract_text(tmp_path, ext)
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass
