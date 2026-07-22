"""
templates.py — All 5 document generators as direct Python functions.
No subprocess. Called directly from app.py.

All 4 XML fixes preserved:
  FIX A: XML reference read as raw text (handled in extractor.py)
  FIX B: sanitize_xml_ref() strips attr values, keeps tag/attr names
  FIX C: XML prompt labels reference as EXACT SCHEMA (in app.py / agent.py)
  FIX D: generate_xml() parses ref_schema to discover real tag/attr names
"""

import os
import json
import uuid
import xml.etree.ElementTree as ET
from xml.dom import minidom
from typing import Any, Dict, Optional


# ── Helpers ───────────────────────────────────────────────────────────────

def _safe(v: Any, maxlen: int = 120) -> str:
    if v is None:
        return ""
    s = str(v)
    s = "".join(c if 32 <= ord(c) < 127 else " " for c in s)
    return s[:maxlen].strip()


def _safe_tag(name: str) -> str:
    """Make a string safe to use as an XML tag name."""
    import re
    s = re.sub(r"[^a-zA-Z0-9_\-.]", "_", str(name))
    if s and s[0].isdigit():
        s = "_" + s
    return s or "Field"


# ── FIX B: sanitize_xml_ref ───────────────────────────────────────────────

def sanitize_xml_ref(xml_text: str) -> str:
    """
    Strips attribute VALUES but keeps attribute NAMES intact.
    Claude sees the real schema shape without noise from actual data values.
    """
    if not xml_text:
        return ""
    import re
    s = "".join(c if 32 <= ord(c) < 128 or c in "\t\n\r" else " " for c in xml_text)
    s = re.sub(r'="[^"]*"', '="..."', s)
    s = re.sub(r"\s{3,}", " ", s)
    return s.strip()[:12000]


# ═══════════════════════════════════════════════════════════════════════════
# XLSX
# ═══════════════════════════════════════════════════════════════════════════

def _xlsx_extract_theme(grounding_path: str):
    """
    Extract header fill colour and alt-row fill colour from the first sheet
    of the grounding XLSX.  Returns (hdr_hex, alt_hex) or (None, None) on failure.
    Only reads cells from the first 6 rows; opens without read_only so styles
    are available.  Silently ignores any error.
    """
    try:
        from openpyxl import load_workbook
        twb = load_workbook(grounding_path, data_only=True)
        tws = twb.worksheets[0]
        hdr_hex = None
        alt_hex = None
        for row in tws.iter_rows(min_row=1, max_row=6):
            for cell in row:
                fill = cell.fill
                if fill and fill.fill_type == "solid":
                    fg = fill.fgColor
                    if fg and fg.type == "rgb":
                        rgb = fg.rgb
                        # Strip leading alpha byte if present (e.g. FF1F497D -> 1F497D)
                        if len(rgb) == 8:
                            rgb = rgb[2:]
                        if len(rgb) == 6 and rgb.upper() not in ("000000", "FFFFFF"):
                            if hdr_hex is None:
                                hdr_hex = rgb.upper()
                            elif rgb.upper() != hdr_hex:
                                alt_hex = rgb.upper()
                                break
            if hdr_hex and alt_hex:
                break
        twb.close()
        return hdr_hex, alt_hex
    except Exception:
        return None, None


def generate_xlsx(plan: Dict, output_path: str, grounding_path: str = "") -> str:
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Border, Side, Alignment
    from openpyxl.chart import BarChart, Reference
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    wb.remove(wb.active)

    # ── Try to pull colours from the grounding template ──────────────────
    hdr_hex, alt_hex = None, None
    if grounding_path and grounding_path.lower().endswith((".xlsx", ".xls")):
        hdr_hex, alt_hex = _xlsx_extract_theme(grounding_path)
        if hdr_hex:
            print(f"   Using XLSX template colours: hdr=#{hdr_hex} alt=#{alt_hex or 'default'}")

    HDR_FILL = PatternFill("solid", fgColor=hdr_hex or "1F497D")
    ALT_FILL = PatternFill("solid", fgColor=alt_hex or "DCE6F1")
    HDR_FONT = Font(bold=True, color="FFFFFF", size=11)
    TTL_FONT = Font(bold=True, size=14, color=hdr_hex or "1F497D")
    THIN     = Side(style="thin", color="B8CCE4")
    BORDER   = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
    CENTER   = Alignment(horizontal="center", vertical="center")

    for sheet in plan.get("sheets", []):
        ws      = wb.create_sheet(title=_safe(sheet.get("name", "Sheet"))[:31])
        headers = [_safe(h) for h in sheet.get("headers", [])]
        rows    = sheet.get("rows", [])
        n_cols  = max(len(headers), 1)

        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=n_cols)
        tc = ws.cell(1, 1, _safe(plan.get("title", "Report")))
        tc.font = TTL_FONT
        tc.alignment = CENTER
        ws.row_dimensions[1].height = 28

        for c, h in enumerate(headers, 1):
            cell = ws.cell(2, c, h)
            cell.fill = HDR_FILL
            cell.font = HDR_FONT
            cell.border = BORDER
            cell.alignment = CENTER
        ws.freeze_panes = "A3"

        for r_idx, row in enumerate(rows):
            fill = ALT_FILL if r_idx % 2 == 0 else None
            for c_idx, val in enumerate(row):
                cell = ws.cell(r_idx + 3, c_idx + 1)
                try:
                    sv = str(val)
                    cell.value = float(sv) if sv.replace(".", "", 1).replace("-", "", 1).isdigit() else sv
                except Exception:
                    cell.value = _safe(val)
                cell.border = BORDER
                if fill:
                    cell.fill = fill

        for c_idx, h in enumerate(headers, 1):
            col_vals = [h] + [_safe(r[c_idx - 1]) if c_idx - 1 < len(r) else "" for r in rows]
            width = min(max((len(v) for v in col_vals), default=8) + 4, 50)
            ws.column_dimensions[get_column_letter(c_idx)].width = width

    # Chart on first sheet
    try:
        sheets  = plan.get("sheets", [])
        c_idx   = int(plan.get("chart_sheet", 0))
        if c_idx < len(wb.sheetnames) and c_idx < len(sheets):
            ws_c   = wb[wb.sheetnames[c_idx]]
            rows_c = sheets[c_idx].get("rows", [])
            n      = len(rows_c)
            if n >= 2:
                col   = int(plan.get("chart_data_col", 1)) + 1
                chart = BarChart()
                chart.title = _safe(plan.get("title", "Chart"))[:50]
                chart.style = 10
                data_ref = Reference(ws_c, min_col=col, min_row=2, max_row=n + 2)
                cats_ref = Reference(ws_c, min_col=1, min_row=3, max_row=n + 2)
                chart.add_data(data_ref, titles_from_data=True)
                chart.set_categories(cats_ref)
                chart.width = 18
                chart.height = 12
                ws_c.add_chart(chart, "A" + str(n + 6))
    except Exception as e:
        print(f"   Chart skipped: {e}")

    wb.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# DOCX
# ═══════════════════════════════════════════════════════════════════════════

def generate_docx(plan: Dict, output_path: str, grounding_path: str = "") -> str:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    def set_cell_bg(cell, hex_color: str):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    # ── Try to use grounding DOCX as style template ───────────────────────
    doc = None
    if grounding_path and grounding_path.lower().endswith((".docx", ".doc")):
        try:
            doc = Document(grounding_path)
            # Clear all body content while keeping styles, themes and page layout.
            # The body must keep its <w:sectPr> (section/margin properties).
            body = doc.element.body
            ns   = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            for child in list(body):
                tag = child.tag.split("}")[-1]
                if tag != "sectPr":          # preserve page layout
                    body.remove(child)
            # DOCX spec requires at least one paragraph before sectPr
            empty_p = OxmlElement("w:p")
            sect_pr = body.find(f"{{{ns}}}sectPr")
            if sect_pr is not None:
                body.insert(list(body).index(sect_pr), empty_p)
            else:
                body.append(empty_p)
            print(f"   Using DOCX template: {os.path.basename(grounding_path)}")
        except Exception as e:
            print(f"   DOCX template load failed ({e}), using default style")
            doc = None

    if doc is None:
        doc = Document()
        for sec in doc.sections:
            sec.top_margin    = Inches(1)
            sec.bottom_margin = Inches(1)
            sec.left_margin   = Inches(1.2)
            sec.right_margin  = Inches(1.2)

    p  = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r  = p.add_run(_safe(plan.get("title", "Document"))[:120])
    r.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    p.paragraph_format.space_after = Pt(8)

    if plan.get("subtitle"):
        s  = doc.add_paragraph()
        s.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = s.add_run(_safe(plan["subtitle"])[:120])
        sr.font.size = Pt(13)
        sr.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)
        s.paragraph_format.space_after = Pt(14)

    for section in plan.get("sections", []):
        h = doc.add_heading(_safe(section.get("heading", ""))[:100], level=int(section.get("level", 1)))
        h.paragraph_format.space_before = Pt(10)
        h.paragraph_format.space_after  = Pt(4)

        for para in section.get("paragraphs", []):
            p = doc.add_paragraph(_safe(para, 800))
            p.paragraph_format.space_after = Pt(6)

        for b in section.get("bullets", []):
            doc.add_paragraph(_safe(b, 300), style="List Bullet")

        tbl = section.get("table")
        if tbl:
            hdrs = [_safe(h) for h in tbl.get("headers", [])]
            rows = tbl.get("rows", [])
            if hdrs:
                t  = doc.add_table(rows=1 + len(rows), cols=len(hdrs))
                t.style = "Table Grid"
                hc = t.rows[0].cells
                for i, h in enumerate(hdrs):
                    hc[i].text = h
                    set_cell_bg(hc[i], "1F497D")
                    if hc[i].paragraphs[0].runs:
                        hc[i].paragraphs[0].runs[0].bold = True
                        hc[i].paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                for ri, row in enumerate(rows):
                    bg = "DCE6F1" if ri % 2 == 0 else "FFFFFF"
                    for ci, val in enumerate(row):
                        if ci < len(t.rows[ri + 1].cells):
                            c = t.rows[ri + 1].cells[ci]
                            c.text = _safe(val, 120)
                            set_cell_bg(c, bg)
                doc.add_paragraph()

    doc.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# PDF
# ═══════════════════════════════════════════════════════════════════════════

def generate_pdf(plan: Dict, output_path: str) -> str:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.lib.enums import TA_CENTER

    doc   = SimpleDocTemplate(output_path, pagesize=A4,
                               leftMargin=inch, rightMargin=inch,
                               topMargin=inch, bottomMargin=inch)
    story = []
    styles = getSampleStyleSheet()
    HDR    = colors.HexColor("#1F497D")
    ALT    = colors.HexColor("#DCE6F1")

    ts = ParagraphStyle("T2", parent=styles["Title"],   fontSize=22, textColor=HDR, spaceAfter=6, alignment=TA_CENTER)
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=15, textColor=HDR, spaceAfter=4)
    bd = ParagraphStyle("B2", parent=styles["Normal"],  fontSize=10, spaceAfter=6, leading=14)

    story.append(Paragraph(_safe(plan.get("title", "Report"), 120), ts))
    if plan.get("subtitle"):
        story.append(Paragraph(_safe(plan["subtitle"], 120),
            ParagraphStyle("S", parent=styles["Normal"], fontSize=12,
                           textColor=colors.HexColor("#44729C"), alignment=TA_CENTER, spaceAfter=10)))
    story.append(HRFlowable(width="100%", thickness=2, color=HDR))
    story.append(Spacer(1, 0.15 * inch))

    for sec in plan.get("sections", []):
        if sec.get("heading"):
            story.append(Paragraph(_safe(sec["heading"], 100), h1))
            story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#B8CCE4")))
            story.append(Spacer(1, 0.06 * inch))
        for para in sec.get("paragraphs", []):
            story.append(Paragraph(_safe(para, 800), bd))
        tbl = sec.get("table")
        if tbl:
            hdrs = [_safe(h) for h in tbl.get("headers", [])]
            rows = tbl.get("rows", [])
            if hdrs:
                td   = [[_safe(h) for h in hdrs]] + [[_safe(v, 80) for v in row] for row in rows]
                cw   = (A4[0] - 2 * inch) / max(len(hdrs), 1)
                t    = Table(td, colWidths=[cw] * len(hdrs), repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND",   (0, 0), (-1, 0),  HDR),
                    ("TEXTCOLOR",    (0, 0), (-1, 0),  colors.white),
                    ("FONTNAME",     (0, 0), (-1, 0),  "Helvetica-Bold"),
                    ("FONTSIZE",     (0, 0), (-1, 0),  10),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, ALT]),
                    ("GRID",         (0, 0), (-1, -1), 0.5, colors.HexColor("#B8CCE4")),
                    ("TOPPADDING",   (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.12 * inch))
        story.append(Spacer(1, 0.08 * inch))

    page_num = [0]

    def footer(canvas, doc):
        page_num[0] += 1
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.grey)
        canvas.drawRightString(A4[0] - inch, 0.5 * inch, f"Page {page_num[0]}")
        canvas.restoreState()

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════

def _pptx_get_accent_color(prs):
    """
    Extract the first non-black/non-white accent colour from the slide master's
    theme colour scheme.  Returns RGBColor or None.
    """
    try:
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        clrScheme = prs.slide_master._element.find(f".//{{{ns_a}}}clrScheme")
        if clrScheme is None:
            return None
        for clr_el in clrScheme:
            for sub in clr_el:
                tag = sub.tag.split("}")[-1]
                if tag == "srgbClr":
                    val = sub.get("val", "")
                    if len(val) == 6 and val.upper() not in ("000000", "FFFFFF", "FFFFFE"):
                        r, g, b = int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16)
                        return RGBColor(r, g, b)
    except Exception:
        pass
    return None


def generate_pptx(plan: Dict, output_path: str, grounding_path: str = "") -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # ── Try to load grounding PPTX as design template ─────────────────────
    prs = None
    used_template = False
    if grounding_path and grounding_path.lower().endswith(".pptx"):
        try:
            prs = Presentation(grounding_path)
            # Remove all existing content slides while preserving slide master + layouts
            sldIdLst = prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                rId = sldId.get(qn("r:id"))
                prs.part.drop_rel(rId)
                sldIdLst.remove(sldId)
            used_template = True
            print(f"   Using PPTX template: {os.path.basename(grounding_path)}")
        except Exception as e:
            print(f"   PPTX template load failed ({e}), using default style")
            prs = None

    if prs is None:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    # ── Colour palette — prefer extracted template accent, fall back to defaults ──
    _accent = _pptx_get_accent_color(prs) if used_template else None
    HDR = _accent or RGBColor(0x1F, 0x49, 0x7D)
    ACC = RGBColor(0x44, 0x72, 0xC4)
    WHT = RGBColor(0xFF, 0xFF, 0xFF)
    BG  = RGBColor(0xF2, 0xF7, 0xFF)

    def set_bg(slide, rgb):
        # When using a template, skip overriding background so the master design shows
        if used_template:
            return
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = rgb

    def add_text(tf, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        tf.text = ""
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = _safe(text, 140)
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    # Use blank layout — pick the one least likely to inject placeholder content
    _layouts = prs.slide_layouts
    blank = _layouts[min(6, len(_layouts) - 1)]

    for sl in plan.get("slides", []):
        t = sl.get("type", "bullets")

        if t == "title":
            slide = prs.slides.add_slide(blank)
            set_bg(slide, BG)
            b = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
            add_text(b.text_frame, sl.get("title", ""), 40, True, HDR, PP_ALIGN.CENTER)
            if sl.get("subtitle"):
                b2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(0.8))
                add_text(b2.text_frame, sl["subtitle"], 20, False, ACC, PP_ALIGN.CENTER)

        elif t == "bullets":
            slide = prs.slides.add_slide(blank)
            set_bg(slide, WHT)
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = HDR
            bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.9))
            add_text(tb.text_frame, sl.get("title", ""), 22, True, WHT)
            bx = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(5.5))
            tf = bx.text_frame
            tf.word_wrap = True
            tf.text = ""
            for i, b in enumerate(sl.get("bullets", [])):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = f"  \u2022  {_safe(b, 120)}"
                p.space_after = Pt(8)
                if p.runs:
                    p.runs[0].font.size = Pt(16)

        elif t == "table":
            slide = prs.slides.add_slide(blank)
            set_bg(slide, WHT)
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = HDR
            bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.9))
            add_text(tb.text_frame, sl.get("title", ""), 22, True, WHT)
            hdrs = sl.get("headers", [])
            rows = sl.get("rows", [])
            if hdrs and rows:
                nc  = len(hdrs)
                nr  = len(rows)
                tbl = slide.shapes.add_table(
                    nr + 1, nc,
                    Inches(0.5), Inches(1.4),
                    Inches(12.3), Inches(min(nr * 0.5 + 0.5, 5.5))
                ).table
                for c, h in enumerate(hdrs):
                    cell = tbl.cell(0, c)
                    cell.text = _safe(h, 50)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = HDR
                    if cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHT
                        cell.text_frame.paragraphs[0].runs[0].font.bold = True
                for r, row in enumerate(rows):
                    bg = RGBColor(0xDC, 0xE6, 0xF1) if r % 2 == 0 else WHT
                    for c, val in enumerate(row):
                        if c < nc:
                            cell = tbl.cell(r + 1, c)
                            cell.text = _safe(val, 80)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = bg

        elif t == "chart":
            slide = prs.slides.add_slide(blank)
            set_bg(slide, WHT)
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = HDR
            bar.line.fill.background()
            tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.9))
            add_text(tb.text_frame, sl.get("title", ""), 22, True, WHT)
            cats = sl.get("categories", ["A", "B", "C"])
            vals = sl.get("values", [1, 2, 3])
            try:
                cd = ChartData()
                cd.categories = cats
                cd.add_series(_safe(sl.get("series_name", "Series"), 40), vals)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5), cd
                ).chart
                chart.has_title = True
                chart.chart_title.text_frame.text = _safe(sl.get("title", ""), 80)
                chart.series[0].format.fill.solid()
                chart.series[0].format.fill.fore_color.rgb = ACC
            except Exception as e:
                print(f"   Chart error: {e}")

    prs.save(output_path)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# XML  — FIX D: reference-driven
# ═══════════════════════════════════════════════════════════════════════════

def generate_xml(plan: Dict, output_path: str, ref_schema: str = "") -> str:
    """
    FIX D: Parses ref_schema at runtime to discover real root tag,
    entity tag names, and field names. Falls back to generic structure
    when no ref_schema is provided.
    """

    def prettify(elem) -> str:
        rough    = ET.tostring(elem, encoding="unicode")
        reparsed = minidom.parseString(rough)
        return reparsed.toprettyxml(indent="  ", encoding=None)

    # ── Parse reference schema ─────────────────────────────────────────────
    root_tag   = "Document"
    ns_uri     = None
    entity_map: Dict[str, list] = {}   # entity_name -> [field_names]
    entity_tag = "Entity"

    if ref_schema and ref_schema.strip().startswith("<"):
        try:
            schema_text = ref_schema
            if schema_text.lstrip().startswith("<?xml"):
                schema_text = schema_text[schema_text.index("?>") + 2:].strip()

            ref_root = ET.fromstring(schema_text)

            raw_tag = ref_root.tag
            if "}" in raw_tag:
                ns_uri   = raw_tag.split("}")[0][1:]
                root_tag = raw_tag.split("}")[1]
            else:
                root_tag = raw_tag

            for child in ref_root:
                ctag_full = child.tag
                ctag      = ctag_full.split("}")[1] if "}" in ctag_full else ctag_full
                entity_tag = ctag

                ent_name = (
                    child.get("name") or child.get("id") or
                    child.get("code") or ctag
                )

                field_names = []
                for subchild in child.iter():
                    stag  = subchild.tag.split("}")[1] if "}" in subchild.tag else subchild.tag
                    fname = (
                        subchild.get("name") or subchild.get("fieldName") or
                        subchild.get("id")
                    )
                    if fname and stag.lower() in ("field", "column", "attribute", "property", "fielddef"):
                        field_names.append(fname)

                entity_map[ent_name] = field_names

            print(f"   Parsed ref schema: root={root_tag}, entities={list(entity_map.keys())}")

        except Exception as parse_err:
            print(f"   ⚠ Could not parse ref_schema ({parse_err}), using generic structure")

    # ── Build output XML ───────────────────────────────────────────────────
    root_el = ET.Element(root_tag)
    if ns_uri:
        root_el.set("xmlns", ns_uri)
    root_el.set("version", "1.0")

    meta = ET.SubElement(root_el, "Metadata")
    ET.SubElement(meta, "Title").text       = _safe(plan.get("title", "Document"))
    ET.SubElement(meta, "GeneratedBy").text = "Claude Agent - ProjectZen"
    ET.SubElement(meta, "Model").text       = "claude-opus-4-7"

    for sheet in plan.get("sheets", []):
        sname   = _safe(sheet.get("name", "Entity")).replace(" ", "")
        headers = sheet.get("headers", [])
        rows    = sheet.get("rows", [])

        # Match sheet name to a reference entity
        ref_fields: list = []
        matched_ent = sname
        for ent_name, fields in entity_map.items():
            if (ent_name.lower() == sname.lower() or
                    ent_name.lower() in sname.lower() or
                    sname.lower() in ent_name.lower()):
                ref_fields  = fields
                matched_ent = ent_name
                break

        ent_el = ET.SubElement(root_el, entity_tag)
        ent_el.set("name",  matched_ent)
        ent_el.set("label", _safe(sheet.get("name", matched_ent)))

        use_fields = ref_fields if ref_fields else headers
        if use_fields:
            fields_el = ET.SubElement(ent_el, "Fields")
            for fname in use_fields:
                f_el = ET.SubElement(fields_el, "Field")
                f_el.set("name",      _safe_tag(fname))
                f_el.set("label",     _safe(fname))
                f_el.set("type",      "String")
                f_el.set("maxLength", "255")

        if rows:
            records_el = ET.SubElement(ent_el, "Records")
            for row in rows:
                record = ET.SubElement(records_el, "Record")
                for i, val in enumerate(row):
                    if i < len(headers):
                        tag_name = _safe_tag(headers[i])
                        ET.SubElement(record, tag_name).text = _safe(val)

    for sec in plan.get("sections", []):
        cc = ET.SubElement(root_el, "Section")
        cc.set("heading", _safe(sec.get("heading", "General")))
        for para in sec.get("paragraphs", []):
            ET.SubElement(cc, "Paragraph").text = _safe(para, 500)

    # ── Write output ───────────────────────────────────────────────────────
    xml_str   = prettify(root_el)
    xml_lines = xml_str.split("\n")
    if xml_lines and xml_lines[0].startswith("<?xml"):
        xml_lines = xml_lines[1:]

    declaration = '<?xml version="1.0" encoding="UTF-8"?>'
    output      = declaration + "\n" + "\n".join(xml_lines)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(output)

    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# CSF XML — SAP SuccessFactors Country-Specific Fields
# Dedicated generator: reads CSF sheets from Excel directly.
# Produces exact <country-specific-fields> structure matching SF DTD.
# ═══════════════════════════════════════════════════════════════════════════

# Sheet name → HRIS element id + column layout
_CSF_SHEET_META = {
    "CSF Personal Info (Global Info)": {
        "hris_element":  "globalInfo",
        "header_row":    3,      # 0-indexed row of the column headers
        "country_col":   2,
        "field_id_col":  3,
        "label_col":     4,
        "type_col":      5,
        "maxlen_col":    6,
        "visibility_col":7,
        "required_col":  9,
        "picklist_col":  10,
    },
    "CSF Addresses": {
        "hris_element":  "homeAddress",
        "header_row":    3,
        "country_col":   2,
        "field_id_col":  3,
        "label_col":     4,
        "type_col":      5,
        "maxlen_col":    6,
        "visibility_col":7,
        "required_col":  9,
        "picklist_col":  10,
    },
    "CSF Dependents": {
        # Has an extra 'Classification' column at index 3 — shifts field cols right by 1
        "hris_element":  "dependents",
        "header_row":    3,
        "country_col":   2,
        "field_id_col":  4,
        "label_col":     5,
        "type_col":      6,
        "maxlen_col":    7,
        "visibility_col":None,   # no visibility column in this sheet
        "required_col":  10,
        "picklist_col":  11,
    },
    "CSF Job Info": {
        "hris_element":  "jobInfo",
        "header_row":    3,
        "country_col":   2,
        "field_id_col":  3,
        "label_col":     4,
        "type_col":      5,
        "maxlen_col":    6,
        "visibility_col":7,
        "required_col":  9,
        "picklist_col":  10,
    },
}

# SAP SF DOCTYPE declaration
_CSF_DOCTYPE = (
    '<?xml version="1.0" encoding="UTF-8"?>\n'
    '<!DOCTYPE country-specific-fields PUBLIC\n'
    ' "-//SuccessFactors, Inc.//DTD Country Specific Field Configuration//EN"\n'
    ' "http://svn/viewvc/svn/V4/trunk/src/com/sf/dtd/country-specific-fields.dtd?view=co">'
)


def _csf_get(row: tuple, col) -> Optional[str]:
    """Safely get a cell value from a row tuple."""
    if col is None or col >= len(row):
        return None
    v = row[col]
    if v is None:
        return None
    s = str(v).strip()
    return s if s and s.lower() != "none" else None


def _csf_country_code(val: str) -> str:
    """
    Extract 3-letter ISO-3166-1 alpha-3 code from a country cell value.

    Handles two formats found in the workbook:
      1. 'Angola (AGO)'     → 'AGO'   (preferred — has code in parentheses)
      2. 'Angola'           → 'AGO'   (fallback — name-to-code lookup)
      3. 'South Korea'      → 'KOR'   (common-name lookup)
    """
    import re

    # ── Primary: code in parentheses ──────────────────────────────────────
    m = re.search(r"\(([A-Z]{2,3})\)", str(val))
    if m:
        return m.group(1)

    # ── Fallback: country name lookup ─────────────────────────────────────
    # Maps common names (and aliases) that appear in SF workbooks
    # to their ISO-3166-1 alpha-3 codes.
    _NAME_TO_CODE = {
        "afghanistan": "AFG", "albania": "ALB", "algeria": "DZA",
        "angola": "AGO", "argentina": "ARG", "armenia": "ARM",
        "australia": "AUS", "austria": "AUT", "azerbaijan": "AZE",
        "bahrain": "BHR", "bangladesh": "BGD", "belarus": "BLR",
        "belgium": "BEL", "bolivia": "BOL", "bosnia": "BIH",
        "botswana": "BWA", "brazil": "BRA", "bulgaria": "BGR",
        "cambodia": "KHM", "cameroon": "CMR", "canada": "CAN",
        "chile": "CHL", "china": "CHN", "colombia": "COL",
        "costa rica": "CRI", "croatia": "HRV", "czech republic": "CZE",
        "czechia": "CZE", "denmark": "DNK", "dominican republic": "DOM",
        "ecuador": "ECU", "egypt": "EGY", "el salvador": "SLV",
        "estonia": "EST", "ethiopia": "ETH", "finland": "FIN",
        "france": "FRA", "georgia": "GEO", "germany": "DEU",
        "ghana": "GHA", "greece": "GRC", "guatemala": "GTM",
        "honduras": "HND", "hong kong": "HKG", "hungary": "HUN",
        "india": "IND", "indonesia": "IDN", "iran": "IRN",
        "iraq": "IRQ", "ireland": "IRL", "israel": "ISR",
        "italy": "ITA", "ivory coast": "CIV", "jamaica": "JAM",
        "japan": "JPN", "jordan": "JOR", "kazakhstan": "KAZ",
        "kenya": "KEN", "kosovo": "XKX", "kuwait": "KWT",
        "latvia": "LVA", "lebanon": "LBN", "lithuania": "LTU",
        "luxembourg": "LUX", "malaysia": "MYS", "mauritius": "MUS",
        "mexico": "MEX", "moldova": "MDA", "morocco": "MAR",
        "mozambique": "MOZ", "myanmar": "MMR", "namibia": "NAM",
        "netherlands": "NLD", "new zealand": "NZL", "nigeria": "NGA",
        "norway": "NOR", "oman": "OMN", "pakistan": "PAK",
        "panama": "PAN", "paraguay": "PRY", "peru": "PER",
        "philippines": "PHL", "poland": "POL", "portugal": "PRT",
        "puerto rico": "PRI", "qatar": "QAT", "romania": "ROU",
        "russia": "RUS", "russian federation": "RUS",
        "saudi arabia": "SAU", "senegal": "SEN", "serbia": "SRB",
        "singapore": "SGP", "slovakia": "SVK", "slovenia": "SVN",
        "south africa": "ZAF", "south korea": "KOR",
        "korea, republic of": "KOR", "republic of korea": "KOR",
        "spain": "ESP", "sri lanka": "LKA", "sweden": "SWE",
        "switzerland": "CHE", "taiwan": "TWN", "tanzania": "TZA",
        "thailand": "THA", "tunisia": "TUN", "turkey": "TUR",
        "turkiye": "TUR", "ukraine": "UKR",
        "united arab emirates": "ARE", "uae": "ARE",
        "united kingdom": "GBR", "uk": "GBR",
        "united states": "USA", "usa": "USA", "us": "USA",
        "uruguay": "URY", "venezuela": "VEN", "vietnam": "VNM",
        "viet nam": "VNM", "zambia": "ZMB", "zimbabwe": "ZWE",
    }

    key = str(val).strip().lower()
    if key in _NAME_TO_CODE:
        return _NAME_TO_CODE[key]

    # ── Last resort: strip non-alpha and take first 3 uppercase chars ──────
    # This is intentionally kept but only reached for truly unrecognised names.
    letters = re.sub(r"[^A-Za-z]", "", str(val))
    return letters.upper()[:3]


def _csf_map_visibility(val: Optional[str]) -> str:
    """Map workbook visibility to SF XML values: both | none | view."""
    if not val:
        return "both"
    v = val.strip().lower()
    if v in ("no", "none", "hide", "hidden"):
        return "none"
    if v in ("view", "view only", "read only", "read"):
        return "view"
    return "both"   # Yes / Edit / both / anything else → both


def _csf_map_required(val: Optional[str]) -> Optional[str]:
    """Return 'true' if required, None otherwise (omit attribute)."""
    if not val:
        return None
    return "true" if val.strip().lower() in ("yes", "true", "mandatory", "y") else None


def generate_csf_xml(excel_path: str, output_path: str) -> str:
    """
    Read all CSF-prefixed sheets from the SuccessFactors configuration
    workbook and produce a country-specific-fields XML file matching the
    SF DTD structure:

        <country-specific-fields>
          <country id="AGO">
            <hris-element id="homeAddress">
              <hris-field id="address1" max-length="256"
                          visibility="both" required="true">
                <label>Care Of</label>
              </hris-field>
              ...
            </hris-element>
            <hris-element id="globalInfo">...</hris-element>
            ...
          </country>
          ...
        </country-specific-fields>
    """
    from openpyxl import load_workbook
    from collections import defaultdict, OrderedDict

    wb = load_workbook(excel_path, read_only=True, data_only=True)

    # country_code → hris_element_id → [field dicts]
    # Use OrderedDict to keep country insertion order (alphabetical by code)
    country_data: Dict[str, Dict[str, list]] = defaultdict(
        lambda: defaultdict(list)
    )

    csf_sheets_found = []

    for sheet_name, meta in _CSF_SHEET_META.items():
        if sheet_name not in wb.sheetnames:
            print(f"   ⚠ Sheet not found: {sheet_name} — skipping")
            continue

        csf_sheets_found.append(sheet_name)
        ws       = wb[sheet_name]
        hris_el  = meta["hris_element"]
        hdr_row  = meta["header_row"]

        for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
            # Skip header rows
            if row_idx <= hdr_row:
                continue

            country_raw = _csf_get(row, meta["country_col"])
            field_id    = _csf_get(row, meta["field_id_col"])

            # Skip blank or header-like rows
            if not country_raw or not field_id:
                continue
            if country_raw.lower().strip(" ") in ("country", "country  "):
                continue

            code = _csf_country_code(country_raw)
            if len(code) < 2:
                continue

            label    = _csf_get(row, meta["label_col"]) or field_id
            ftype    = _csf_get(row, meta["type_col"])   or "STRING"
            maxlen   = _csf_get(row, meta["maxlen_col"]) or "256"
            vis_raw  = _csf_get(row, meta["visibility_col"])
            req_raw  = _csf_get(row, meta["required_col"])
            picklist = _csf_get(row, meta["picklist_col"])

            # Normalise max-length to integer string
            try:
                maxlen = str(int(float(str(maxlen))))
            except (ValueError, TypeError):
                maxlen = "256"

            country_data[code][hris_el].append({
                "id":         field_id,
                "label":      label,
                "type":       ftype,
                "max_length": maxlen,
                "visibility": _csf_map_visibility(vis_raw),
                "required":   _csf_map_required(req_raw),
                "picklist":   picklist,
            })

    wb.close()

    if not csf_sheets_found:
        raise ValueError(
            "No CSF sheets found in the workbook. "
            "Expected sheets named: " + ", ".join(_CSF_SHEET_META.keys())
        )

    print(f"   ✅ CSF extraction: {len(country_data)} countries "
          f"from {len(csf_sheets_found)} sheets")

    # ── Build XML tree ────────────────────────────────────────────────────
    root = ET.Element("country-specific-fields")

    for code in sorted(country_data.keys()):
        country_el = ET.SubElement(root, "country")
        country_el.set("id", code)

        for hris_el_id, fields in country_data[code].items():
            if not fields:
                continue
            hris_el = ET.SubElement(country_el, "hris-element")
            hris_el.set("id", hris_el_id)

            for f in fields:
                hf = ET.SubElement(hris_el, "hris-field")
                hf.set("id",         f["id"])
                hf.set("max-length", f["max_length"])
                hf.set("visibility", f["visibility"])
                if f["required"]:
                    hf.set("required", f["required"])
                if f["picklist"]:
                    hf.set("picklist", f["picklist"])

                label_el = ET.SubElement(hf, "label")
                label_el.text = f["label"]

    # ── Serialise with pretty-print ───────────────────────────────────────
    rough    = ET.tostring(root, encoding="unicode")
    reparsed = minidom.parseString(rough)
    pretty   = reparsed.toprettyxml(indent="  ", encoding=None)

    # Strip minidom's own <?xml?> declaration — we use _CSF_DOCTYPE instead
    xml_lines = pretty.split("\n")
    if xml_lines and xml_lines[0].startswith("<?xml"):
        xml_lines = xml_lines[1:]

    output = _CSF_DOCTYPE + "\n" + "\n".join(xml_lines)

    with open(output_path, "w", encoding="utf-8") as fh:
        fh.write(output)

    print(f"   ✅ CSF XML saved → {output_path}")
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# Dispatcher
# ═══════════════════════════════════════════════════════════════════════════

def generate(
    output_format: str,
    plan: Dict,
    output_path: str,
    ref_schema: str = "",
    excel_path: str = "",        # used by csf_xml route
    grounding_path: str = "",    # path to admin-uploaded reference file for template cloning
) -> str:
    """
    Single entry point. Dispatches to the correct generator.
    Returns output_path on success, raises on failure.

    output_format values:
        xlsx | docx | pdf | pptx | xml | csf_xml
    """
    fmt = output_format.lower().strip()
    if fmt == "xlsx":
        return generate_xlsx(plan, output_path, grounding_path)
    if fmt == "docx":
        return generate_docx(plan, output_path, grounding_path)
    if fmt == "pdf":
        return generate_pdf(plan, output_path)
    if fmt == "pptx":
        return generate_pptx(plan, output_path, grounding_path)
    if fmt == "xml":
        return generate_xml(plan, output_path, ref_schema)
    if fmt == "csf_xml":
        if not excel_path:
            raise ValueError("excel_path is required for csf_xml format")
        return generate_csf_xml(excel_path, output_path)
    raise ValueError(f"Unsupported output format: {output_format}")
